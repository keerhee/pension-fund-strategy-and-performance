# -*- coding: utf-8 -*-
"""기존 강의본 덱과 같은 톤으로 슬라이드를 짓는 최소 도구.
색·폰트·좌표는 W01 강의본에서 실측한 값이다."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x1B, 0x2C, 0x5E)
INK    = RGBColor(0x3B, 0x42, 0x52)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x2E, 0x5B, 0xAA)
GREEN  = RGBColor(0x3F, 0xA3, 0x6F)
ORANGE = RGBColor(0xE8, 0x87, 0x3A)
PANEL_B= RGBColor(0xEA, 0xF1, 0xFB)
PANEL_G= RGBColor(0xE6, 0xF5, 0xF0)
PANEL_N= RGBColor(0xF4, 0xF5, 0xF9)
CODE_BG= RGBColor(0xF7, 0xF8, 0xFA)
FONT   = "Noto Sans CJK KR"
MONO   = "Noto Sans Mono CJK KR"
COURSE = "BAF.60080 · 연기금 운용전략과 성과평가"

def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs

def _tb(s, x, y, w, h, text, size, color, bold=False, font=FONT, align=PP_ALIGN.LEFT,
        line=1.25, space_after=4):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0; tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line   # float = 배수, Pt = 절대 행간
        p.space_after = Pt(space_after)
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def _rect(s, x, y, w, h, fill, radius=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    if radius:
        try: shp.adjustments[0] = radius
        except Exception: pass
    shp.text_frame.text = ""
    return shp

def chrome(s, eyebrow, title, subtitle, page, eyebrow_fill=NAVY):
    """머리말 배지 · 제목 · 부제 · 꼬리말 · 쪽번호"""
    w = 0.34 + 0.14 * len(eyebrow)
    _rect(s, 0.6, 0.28, w, 0.40, eyebrow_fill)
    _tb(s, 0.6, 0.33, w, 0.34, eyebrow, 13, WHITE, True, align=PP_ALIGN.CENTER)
    _tb(s, 8.7, 0.36, 4.0, 0.3, COURSE, 12, MUTED, align=PP_ALIGN.RIGHT)
    _tb(s, 0.6, 0.74, 12.1, 0.72, title, 32, NAVY, True)
    if subtitle:
        _tb(s, 0.6, 1.44, 12.1, 0.42, subtitle, 17, MUTED)
    _tb(s, 0.6, 7.04, 8.0, 0.3, COURSE, 11, MUTED)
    _tb(s, 11.9, 7.04, 0.8, 0.3, str(page), 11, MUTED, align=PP_ALIGN.RIGHT)

def _wrapped_lines(text, w_in, fs):
    """한글 기준 한 줄에 들어가는 글자 수로 줄 수를 추정한다.
    Noto Sans CJK 한글 글자폭 ≈ 폰트크기(pt)/72 inch, 라틴은 그 절반쯤."""
    per = max(8, int(w_in / (fs / 72.0 * 0.97)))
    n = 0
    for ln in text.split("\n"):
        cost = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in ln)
        n += max(1, int(cost / per) + (1 if cost % per else 0))
    return n

def panel_height(w, lines, fs=18, line=1.34, head=True):
    """내용이 다 들어가는 패널 높이(inch)를 돌려준다."""
    inner = w - 0.56 - 0.1
    rows = sum(_wrapped_lines("•  " + t, inner, fs) for t in lines)
    body = rows * (fs * line / 72.0) + len(lines) * (7 / 72.0)
    return round((0.82 if head else 0.30) + body + 0.30, 2)

def panel(s, x, y, w, h, head, lines, tone="blue"):
    if h is None:
        h = panel_height(w, lines, head=bool(head))
    bg   = {"blue": PANEL_B, "green": PANEL_G, "gray": PANEL_N}[tone]
    head_c = {"blue": BLUE, "green": GREEN, "gray": NAVY}[tone]
    _rect(s, x, y, w, h, bg, radius=0.03)
    top = 0.82
    if head:
        _tb(s, x + 0.30, y + 0.22, w - 0.6, 0.40, head, 20, head_c, True)
    else:
        top = 0.30
    body = "\n".join(("•  " + t) if not t.startswith(" ") else t.strip() for t in lines)
    _tb(s, x + 0.30, y + top, w - 0.6, h - top - 0.28, body, 18, INK, line=1.34, space_after=7)
    return y + h

CODE_LEAD = 20          # 코드 행간(pt) — 절대값으로 고정한다

def code_height(code, fs=16, line=None):
    n = len(code.split("\n"))
    lead = CODE_LEAD if fs >= 16 else fs + 4
    return round(0.44 + 0.22 + n * (lead / 72.0) + 0.22, 2)

def codebox(s, x, y, w, h, head, code, tone=ORANGE, fs=16):
    if h is None:
        h = code_height(code, fs=fs)
    _rect(s, x, y, w, 0.46, tone)
    _tb(s, x + 0.22, y + 0.08, w - 0.34, 0.34, head, 14, WHITE, True)
    _rect(s, x, y + 0.46, w, h - 0.46, CODE_BG)
    _tb(s, x + 0.26, y + 0.66, w - 0.52, h - 0.88, code, fs, INK, font=MONO,
        line=Pt(CODE_LEAD if fs >= 16 else fs + 4), space_after=0)
    return y + h

def table_height(headers, rows, widths, fs=18, minrow=0.52):
    """행마다 필요한 높이를 더한 표 전체 높이."""
    h = 0.52
    for row in rows:
        need = 1
        for i, cell in enumerate(row):
            need = max(need, _wrapped_lines(str(cell), widths[i] - 0.34, fs))
        h += max(minrow, need * (fs * 1.24 / 72.0) + 0.20)
    return round(h, 2)

def table(s, x, y, w, headers, rows, widths=None, fs=18, rowh=None):
    """행 높이를 셀 내용에서 자동으로 잡는다 — 두 줄짜리 셀이 다음 행을 덮지 않게."""
    n = len(headers)
    widths = widths or [w / n] * n
    _rect(s, x, y, w, 0.52, NAVY)
    cx = x
    for i, htxt in enumerate(headers):
        _tb(s, cx + 0.16, y + 0.13, widths[i] - 0.32, 0.34, htxt, 17, WHITE, True)
        cx += widths[i]
    yy = y + 0.52
    for j, row in enumerate(rows):
        need = 1
        for i, cell in enumerate(row):
            need = max(need, _wrapped_lines(str(cell), widths[i] - 0.34, fs))
        rh = max(rowh or 0.52, need * (fs * 1.24 / 72.0) + 0.20)
        if j % 2 == 1:
            _rect(s, x, yy, w, rh, PANEL_N)
        cx = x
        for i, cell in enumerate(row):
            _tb(s, cx + 0.16, yy + 0.09, widths[i] - 0.32, rh - 0.16, str(cell), fs, INK,
                line=1.24, space_after=0)
            cx += widths[i]
        yy += rh
    return yy

def band(s, y, text, fill=NAVY, color=WHITE, size=18, x=0.6, w=12.1, h=0.62):
    _rect(s, x, y, w, h, fill)
    _tb(s, x + 0.34, y + 0.16, w - 0.68, h - 0.32, text, size, color, True)
    return y + h

def cover(s, week, title, subtitle, tagline):
    bg = _rect(s, 0, 0, 13.333, 7.5, NAVY)
    _tb(s, 1.0, 1.5, 3.0, 1.4, week, 64, ORANGE, True)
    _rect(s, 1.0, 3.05, 3.3, 0.42, ORANGE)
    _tb(s, 1.0, 3.09, 3.3, 0.36, "실습 데이터 가이드", 15, WHITE, True, align=PP_ALIGN.CENTER)
    tfs = 38
    while _wrapped_lines(title, 11.0, tfs) > 1 and tfs > 27:
        tfs -= 1
    _tb(s, 1.0, 3.75 + (0.10 if tfs < 34 else 0), 11.0, 1.0, title, tfs, WHITE, True)
    _tb(s, 1.0, 4.85, 11.0, 0.5, subtitle, 19, RGBColor(0x8F, 0xA8, 0xD8))
    _tb(s, 1.0, 5.66, 11.0, 0.9, tagline, 16, RGBColor(0xC7, 0xD3, 0xEC), line=1.45)
    _tb(s, 1.0, 6.85, 11.0, 0.3, COURSE + "  ·  2026 가을학기", 13, RGBColor(0x8F, 0xA8, 0xD8))
