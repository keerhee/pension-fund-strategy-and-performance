# -*- coding: utf-8 -*-
"""옛 주차 번호가 그림 안에 박혀 있던 도해 두 장을 다시 그린다.

  W04_M4 슬라이드 15 — MVO 4대 한계 → 해법 매핑 : "Black-Litterman (W5)" → 같은 주 2교시
  W15_케이스 슬라이드 10 — 4기관 TPA 전환 로드맵 : "Week 11-16 · W11~W16" → Week 10–15 · W10~W15

    .venv/bin/python _build/fig_fix_art.py            # PNG만 생성
    .venv/bin/python _build/fig_fix_art.py --embed    # pptx의 그림까지 교체
"""
import os
import sys

import matplotlib
matplotlib.use("Agg")
from matplotlib import pyplot as plt, font_manager as fm
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

NAVY, BLUE, GREEN = "#1b2c5e", "#2e5baa", "#3fa36f"
RED, REDP = "#b3392f", "#fbeceb"
GREENP, BLUEP = "#e8f5ee", "#eaf1fb"
INK, MUTED, HAIR = "#3b4252", "#6b7280", "#d4d9e4"

_FD = os.path.expanduser("~/Library/Fonts")
for _f in ("Pretendard-Regular.otf", "Pretendard-SemiBold.otf", "Pretendard-Bold.otf"):
    _p = os.path.join(_FD, _f)
    if os.path.exists(_p):
        fm.fontManager.addfont(_p)
plt.rcParams["font.family"] = "Pretendard"
plt.rcParams["axes.unicode_minus"] = False

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "_art", "fix")
os.makedirs(OUT, exist_ok=True)


def canvas(w, h):
    fig = plt.figure(figsize=(w / 100, h / 100), dpi=100)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, w); ax.set_ylim(0, h); ax.invert_yaxis(); ax.axis("off")
    fig.patch.set_facecolor("white")
    return fig, ax


def card(ax, x, y, w, h, fc="white", ec=HAIR, lw=1.6, r=12):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle=f"round,pad=0,rounding_size={r}",
                                facecolor=fc, edgecolor=ec, linewidth=lw))


def t(ax, x, y, s, size=14, color=INK, weight="normal", ha="center"):
    ax.text(x, y, s, fontsize=size, color=color, ha=ha, va="center",
            fontweight=weight, linespacing=1.5)


def arrow(ax, x0, y0, x1, y1, color=MUTED, lw=2.2):
    ax.add_patch(FancyArrowPatch((x0, y0), (x1, y1), arrowstyle="-|>",
                                 mutation_scale=16, linewidth=lw, color=color,
                                 shrinkA=0, shrinkB=0))


def save(fig, name):
    p = os.path.join(OUT, name)
    fig.savefig(p, dpi=100, facecolor="white")
    plt.close(fig)
    print("saved", p)
    return p


# ── W04_M4 s15 — MVO의 4대 한계 → 해법 매핑 ──────────────────────
def mvo_limits():
    W, H = 1600, 640
    fig, ax = canvas(W, H)
    t(ax, W / 2, 34, "MVO의 4대 한계 → 해법 매핑 — 비판은 폐기가 아니라 보완으로", 21, NAVY, "bold")

    rows = [
        ("① 추정 오차 증폭", "Error Maximizer (Michaud)",
         "Ledoit-Wolf 수축 · RMT", "공분산의 노이즈를 제거"),
        ("② 코너해 집중", "소수 자산에 ~100%",
         "비중 제약 · Robust (L1/L2)", "불확실성을 페널티로 전환"),
        ("③ 비안정성", "입력 ±0.5%에 20~30%p 요동",
         "Resampled Efficiency", "부트스트랩 평균으로 평활화"),
        ("④ 비현실", "레버리지 · 집중 — 실행 불가",
         "제약 + Black-Litterman (W4 · 2교시)", "균형 기반으로 극단 방지"),
    ]
    M, GAP_X, GAP_Y = 118, 196, 20
    cw = (W - 2 * M - GAP_X) / 2
    y0, ch = 86, 106
    for i, (lh, lb, rh, rb) in enumerate(rows):
        y = y0 + i * (ch + GAP_Y)
        card(ax, M, y, cw, ch, fc=REDP, ec="#e6b8b3", lw=1.4)
        t(ax, M + cw / 2, y + 36, lh, 16, RED, "bold")
        t(ax, M + cw / 2, y + 72, lb, 14, INK)
        rx = M + cw + GAP_X
        card(ax, rx, y, cw, ch, fc=GREENP, ec="#a9d9be", lw=1.4)
        t(ax, rx + cw / 2, y + 36, rh, 16, GREEN, "bold")
        t(ax, rx + cw / 2, y + 72, rb, 14, INK)
        arrow(ax, M + cw + 26, y + ch / 2, rx - 26, y + ch / 2)

    t(ax, W / 2, 606,
      "공통 처방 — “μ · Σ를 정확히 안다”는 과신을 버리는 것 : 입력을 고치거나, 불확실성을 명시하거나",
      16, NAVY, "bold")
    return save(fig, "w04_m4_limits.png")


# ── W15_케이스 s10 — 4기관 5년 TPA 전환 로드맵 ────────────────────
def tpa_roadmap():
    W, H = 1600, 640
    fig, ax = canvas(W, H)
    t(ax, W / 4 + 20, 30, "A. 4기관 현재 → 5년 후 목표", 17, NAVY, "bold")
    t(ax, 3 * W / 4 - 20, 30, "B. 단계적 전환 + 통합 시너지", 17, NAVY, "bold")

    left = [
        ("NPS (1,458조) — 부분 적용",
         ["현재: 자산군별 사일로 4부서 · 시장 영향력",
          "5년 후: 60% TPA · One CIO + Factor 위험",
          "시너지: +1.0~1.5%p/년"]),
        ("KIC (USD 232B) — 학술적 적합",
         ["현재: 부분 TPA (Hoon Lee CIO) · 글로벌 경험",
          "5년 후: 90% TPA 정착 (CPPIB 학습 완료)",
          "시너지: +1.5~2.0%p/년"]),
        ("KSWF (20조, 2026.6) — 최적",
         ["현재: 출범 전 · 학술 자문 위원회 구성",
          "5년 후: 100% TPA 출발 정착 (Norway·GIC)",
          "시너지: +1.8~2.5%p/년"]),
        ("국민성장펀드 (150조) — 부분 적용",
         ["현재: 5년 폐쇄 + 12개 섹터 PM 사일로",
          "5년 후: 부분 TPA · One Fund 부분 도입",
          "시너지: +0.5~1.0%p/년"]),
    ]
    right = [
        ("1년차 (2027) — 인프라 구축",
         ["학술 자문 위원회 (G. Rubin + 학계)",
          "CAIA + Thinking Ahead Institute 협력 · KIC 우선"]),
        ("2년차 (2028) — Factor 도입",
         ["Factor-based 위험 측정 도입",
          "전체 펀드 총수익률 평가 · 파생·오버레이"]),
        ("3년차 (2029) — One Fund culture",
         ["부서장 권한 부분 해체 · 통합 CIO 거버넌스",
          "동적 배분 시범 · 4기관 통합 위험 모니터링"]),
        ("4-5년차 (2030-31) — 완성 + 검증",
         ["월/분기 동적 배분 본격화 · 시장 영향력 최소화",
          "Thinking Ahead Institute 정량 검증 참여"]),
        ("통합 시너지 (Week 10-15)",
         ["+8.3~9.8%p/년 (W10 +3.5 … W15 +0.75)",
          "5년 누적 약 120~150조 실현 (잠식 후)"]),
    ]

    M, MID = 22, 30
    cw = (W - 2 * M - MID) / 2
    ly, lh, lgap = 54, 136, 8
    for i, (head, lines) in enumerate(left):
        y = ly + i * (lh + lgap)
        card(ax, M, y, cw, lh, fc=BLUEP, ec=BLUE, lw=1.6)
        t(ax, M + 24, y + 30, head, 15, BLUE, "bold", ha="left")
        for j, ln in enumerate(lines):
            t(ax, M + 24, y + 62 + 28 * j, ln, 13, INK, ha="left")

    rx = M + cw + MID
    ry, rh, rgap = 56, 96, 15
    for i, (head, lines) in enumerate(right):
        y = ry + i * (rh + rgap)
        card(ax, rx, y, cw, rh, fc=GREENP, ec=GREEN, lw=1.6)
        t(ax, rx + 24, y + 28, head, 15, GREEN, "bold", ha="left")
        for j, ln in enumerate(lines):
            t(ax, rx + 24, y + 58 + 26 * j, ln, 13, INK, ha="left")
    return save(fig, "w15_case_roadmap.png")


SLOTS = [
    ("W04_MVO와블랙리터맨/W04_M4_MVO와공분산추정_강의본.pptx", 15, "w04_m4_limits.png"),
    ("W15_TPA/W15_케이스_KFP_TPA전환과평가체계_IC.pptx", 10, "w15_case_roadmap.png"),
]


def embed():
    from pptx import Presentation
    root = os.path.dirname(HERE)
    for rel, sn, png in SLOTS:
        path = os.path.join(root, rel)
        prs = Presentation(path)
        pics = [sh for sh in prs.slides[sn - 1].shapes if sh.shape_type == 13]
        if len(pics) != 1:
            print("SKIP", rel, sn, len(pics)); continue
        pic = pics[0]
        part = pic.part.related_part(pic._element.blip_rId)
        part._blob = open(os.path.join(OUT, png), "rb").read()
        prs.save(path)
        print("embedded", png, "→", rel, f"s{sn}")


if __name__ == "__main__":
    mvo_limits()
    tpa_roadmap()
    if "--embed" in sys.argv:
        embed()
