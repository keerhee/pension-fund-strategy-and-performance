# -*- coding: utf-8 -*-
"""도해 안에 박혀 있는 낱말만 골라 다시 찍는다 — 옛 주차 번호 교정용.

그림 전체를 다시 그리면 나머지 요소가 미묘하게 달라진다. 바꿀 것이 "W16 → W15"
한 낱말뿐일 때는 그 낱말의 자리만 배경색으로 지우고 같은 글꼴 · 같은 크기로 다시 찍는다.

낱말 상자는 tesseract 로 찾고, 글꼴 · 크기 · 위치는 옛 낱말을 여러 조합으로 렌더해
원본과 가장 닮은 것을 골라 정한다(자기 보정). 그래서 손으로 맞출 값이 없다.

    .venv/bin/python _build/img_word_patch.py           # 미리보기 PNG만
    .venv/bin/python _build/img_word_patch.py --embed   # pptx의 그림까지 교체
"""
import os
import re
import subprocess
import sys

from PIL import Image, ImageDraw, ImageFont

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(HERE, "_art", "wordpatch")
os.makedirs(OUT, exist_ok=True)


def _fonts():
    cand = [os.path.expanduser("~/Library/Fonts/Pretendard-Bold.otf"),
            os.path.expanduser("~/Library/Fonts/Pretendard-SemiBold.otf"),
            os.path.expanduser("~/Library/Fonts/Pretendard-Regular.otf")]
    try:                                   # 덱 도해 일부는 matplotlib 기본 글꼴로 그려졌다
        import matplotlib
        d = os.path.join(os.path.dirname(matplotlib.__file__), "mpl-data", "fonts", "ttf")
        cand += [os.path.join(d, n) for n in ("DejaVuSans-Bold.ttf", "DejaVuSans.ttf")]
    except Exception:
        pass
    return [p for p in cand if os.path.exists(p)]


FONTS = _fonts()

# (pptx 상대경로, 슬라이드, 그림 순번, [(찾을 낱말, 옛 표기, 새 표기), ...])
JOBS = [
    ("W02_자산배분과CAPM/W02_자산배분과CAPM_강의본.pptx", 21, 0,
     [("W13", "W13의", "W12의")]),
    ("W03_연기금모델과CMA/W03_연기금모델과CMA_강의본.pptx", 33, 0,
     [("W16", "W16", "W15")]),
    ("W04_MVO와블랙리터맨/W04_M4_MVO와공분산추정_강의본.pptx", 41, 0,
     [("W16", "W16", "W15")]),
    ("W07_LDI와GBI/W07_M9_케이스1_퇴직연금디폴트옵션도입_IC.pptx", 7, 0,
     [((344, 233, 17, 21), "7)", "6)"), ((488, 811, 13, 16), "7)", "6)", "Bold")]),
    ("W07_LDI와GBI/W07_M9_케이스1_퇴직연금디폴트옵션도입_IC.pptx", 10, 0,
     [((465, 258, 43, 18), "7-9의", "6-7의"), ((1129, 817, 26, 18), "10)", "9)", "Bold")]),
]

_norm = lambda s: re.sub(r"[^0-9A-Za-z가-힣]", "", s)


def word_boxes(png_path, word):
    """tesseract TSV 에서 낱말과 겹치는 상자를 왼쪽·위 순서로 모두 돌려준다."""
    tsv = subprocess.run(["tesseract", png_path, "-", "-l", "kor+eng", "--psm", "11", "tsv"],
                         capture_output=True, text=True).stdout
    want = _norm(word)
    hits = []
    for line in tsv.splitlines()[1:]:
        c = line.split("\t")
        if len(c) < 12:
            continue
        tok = _norm(c[11])
        if not tok:
            continue
        if tok == want or (len(want) >= 3 and tok.startswith(want)):
            hits.append((int(c[6]), int(c[7]), int(c[8]), int(c[9])))
    hits.sort(key=lambda b: (b[1], b[0]))
    return hits


def flatten(img):
    """투명 배경은 흰 바탕에 얹어 본다 — 글자 모양 비교용."""
    if img.mode == "RGBA":
        bg = Image.new("RGB", img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[3])
        return bg
    return img.convert("RGB")


def ink_mask(img):
    return flatten(img).convert("L").point(lambda v: 255 if v < 170 else 0)


def sample_colors(img, box):
    """글자색은 상자 안에서 가장 어두운 쪽, 배경색은 상자 둘레에서 가장 흔한 쪽."""
    l, t, w, h = box
    rgba = img.convert("RGBA")
    crop = rgba.crop((l, t, l + w, t + h))
    px = [p for p in crop.getdata() if p[3] > 128]
    if not px:
        px = list(crop.getdata())
    lum = sorted(((0.299 * p[0] + 0.587 * p[1] + 0.114 * p[2], p) for p in px),
                 key=lambda x: x[0])
    color = lum[max(0, len(lum) // 20)][1]
    ring = []
    for x in range(max(0, l - 6), min(img.width, l + w + 6)):
        for y in (max(0, t - 7), min(img.height - 1, t + h + 6)):
            ring.append(rgba.getpixel((x, y)))
    bg = max(set(ring), key=ring.count) if ring else (255, 255, 255, 255)
    if img.mode != "RGBA":
        color, bg = color[:3], bg[:3]
    return color, bg


def fit(img, box, old_text, color, bg, font_hint=None):
    l, t, w, h = box
    pad = 10
    target = ink_mask(img.crop((l - pad, t - pad, l + w + pad, t + h + pad)))
    tw, th = target.size
    best = None
    fonts = [f for f in FONTS if not font_hint or font_hint in os.path.basename(f)] or FONTS
    for fp in fonts:
        for size in range(max(8, h - 6), h + 16):
            f = ImageFont.truetype(fp, size)
            for dx in range(-5, 6):
                for dy in range(-9, 10):
                    im = Image.new(img.mode, (tw, th), bg)
                    ImageDraw.Draw(im).text((pad + dx, pad + dy), old_text, font=f,
                                            fill=color, anchor="lt")
                    m = ink_mask(im)
                    diff = sum(1 for a, b in zip(target.getdata(), m.getdata()) if a != b)
                    if best is None or diff < best[0]:
                        best = (diff, fp, size, dx, dy)
    return best


def patch_one(img, box, old_text, new_text, font_hint=None):
    color, bg = sample_colors(img, box)
    diff, fp, size, dx, dy = fit(img, box, old_text, color, bg, font_hint)
    l, t, w, h = box
    pad = 10
    f = ImageFont.truetype(fp, size)
    d = ImageDraw.Draw(img)
    ow = d.textlength(old_text, font=f)
    d.rectangle([l + dx - 2, t - pad + 3, l + dx + ow + 3, t + h + pad - 3], fill=bg)
    d.text((l + dx, t + dy), new_text, font=f, fill=color, anchor="lt")
    return f"{old_text}→{new_text} · {os.path.basename(fp)} {size}px diff={diff}"


def main():
    from pptx import Presentation
    do_embed = "--embed" in sys.argv
    for rel, sn, idx, patches in JOBS:
        path = os.path.join(ROOT, rel)
        prs = Presentation(path)
        pics = [sh for sh in prs.slides[sn - 1].shapes if sh.shape_type == 13]
        if idx >= len(pics):
            print("SKIP", rel, sn); continue
        pic = pics[idx]
        base = f"{os.path.basename(rel)[:-5]}_s{sn}"
        src = os.path.join(OUT, base + "_src.png")
        open(src, "wb").write(pic.image.blob)
        img = Image.open(src)
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGBA" if "transparency" in img.info else "RGB")
        print(rel, f"s{sn}")
        used = {}
        ok = True
        for entry in patches:
            find, old_text, new_text = entry[0], entry[1], entry[2]
            hint = entry[3] if len(entry) > 3 else None
            if isinstance(find, tuple):
                box = find
            else:
                boxes = word_boxes(src, find)
                k = used.get(find, 0)
                if len(boxes) <= k:
                    print("  낱말을 찾지 못했다:", find); ok = False; continue
                used[find] = k + 1
                box = boxes[k]
            print("  " + patch_one(img, box, old_text, new_text, hint))
        dst = os.path.join(OUT, base + "_new.png")
        img.save(dst)
        if ok and do_embed:
            part = pic.part.related_part(pic._element.blip_rId)
            part._blob = open(dst, "rb").read()
            prs.save(path)
            print("  embedded")


if __name__ == "__main__":
    main()
