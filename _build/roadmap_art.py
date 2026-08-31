# -*- coding: utf-8 -*-
"""주차별 강의본의 로드맵 도해 재생성 — 새 16주 배치표 기준.

원본 도해는 옛 모듈 번호(M1~M16)를 주차 번호처럼 달고 있었다.
새 배치표는 M4+M5=W04 · M7+SS2=W06 · M8+M9=W07 · 8주와 16주는 시험이므로
    LDI·GBI=W07 · 위험관리=W09 · 주식퀀트=W10 · 채권=W11 · 팩터=W12
    매크로/CTA=W13 · 대체투자=W14 · TPA=W15 · 캡스톤·기말=W16
이다. 이 스크립트가 그 번호로 도해를 다시 그리고 pptx 안의 그림을 교체한다.

    .venv/bin/python _build/roadmap_art.py          # PNG만 생성
    .venv/bin/python _build/roadmap_art.py --embed  # pptx의 그림까지 교체

스타일은 _build/w01_art.py(카드 · 칩 · 화살표 · 팔레트)를 그대로 따른다.
"""
import os
import sys

import matplotlib
matplotlib.use("Agg")
from matplotlib import pyplot as plt, font_manager as fm
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

NAVY, BLUE, GREEN = "#1b2c5e", "#2e5baa", "#3fa36f"
ORANGE = "#e2570f"
INK, MUTED = "#3b4252", "#6b7280"
PANEL, BLUEP, GREENP = "#f4f5f9", "#eaf1fb", "#e6f5f0"
HAIR = "#d4d9e4"

_FD = os.path.expanduser("~/Library/Fonts")
for _f in ("Pretendard-Regular.otf", "Pretendard-SemiBold.otf", "Pretendard-Bold.otf"):
    _p = os.path.join(_FD, _f)
    if os.path.exists(_p):
        fm.fontManager.addfont(_p)
plt.rcParams["font.family"] = "Pretendard"
plt.rcParams["axes.unicode_minus"] = False

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "_art", "roadmap")
os.makedirs(OUT, exist_ok=True)


def canvas(w, h):
    fig = plt.figure(figsize=(w / 100, h / 100), dpi=100)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, w); ax.set_ylim(0, h); ax.invert_yaxis(); ax.axis("off")
    fig.patch.set_facecolor("white")
    return fig, ax


def card(ax, x, y, w, h, fc="white", ec=HAIR, lw=1.6, r=12):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle=f"round,pad=0,rounding_size={r}",
                                facecolor=fc, edgecolor=ec, linewidth=lw))


def t(ax, x, y, s, size=14, color=INK, weight="normal", ha="center", va="center"):
    ax.text(x, y, s, fontsize=size, color=color, ha=ha, va=va,
            fontweight=weight, linespacing=1.5)


def arrow(ax, x0, y0, x1, y1, color=MUTED, lw=2.2):
    ax.add_patch(FancyArrowPatch((x0, y0), (x1, y1), arrowstyle="-|>",
                                 mutation_scale=16, linewidth=lw, color=color,
                                 shrinkA=0, shrinkB=0))


def save(fig, name):
    p = os.path.join(OUT, name)
    fig.savefig(p, dpi=100, facecolor="white")
    plt.close(fig)
    print("saved", p)
    return p


# ── 1) 다섯 카드 사슬 (1600×640) — W02~W05 강의본 ────────────────
def chain(name, title, steps, banner):
    """steps = [(칩, 제목, [줄1, 줄2], 종류)] · 종류 = today | grey | blue | green"""
    W, H = 1600, 640
    fig, ax = canvas(W, H)
    t(ax, W / 2, 34, title, 21, NAVY, "bold")

    M, GAP = 24, 42
    n = len(steps)
    bw = (W - 2 * M - (n - 1) * GAP) / n
    by, bh = 130, 260
    for i, (chip, head, lines, kind) in enumerate(steps):
        x = M + i * (bw + GAP)
        if kind == "today":
            fc, ec, lw, chipc, tc = "white", NAVY, 2.6, NAVY, NAVY
        elif kind == "blue":
            fc, ec, lw, chipc, tc = BLUEP, BLUE, 1.6, BLUE, BLUE
        elif kind == "green":
            fc, ec, lw, chipc, tc = GREENP, GREEN, 1.6, GREEN, GREEN
        else:
            fc, ec, lw, chipc, tc = PANEL, HAIR, 1.6, NAVY, NAVY
        card(ax, x, by, bw, bh, fc=fc, ec=ec, lw=lw)
        cx = x + bw / 2
        chw, chh = min(bw * 0.72, 210), 44
        card(ax, cx - chw / 2, by - 22, chw, chh, fc=chipc, ec=chipc, r=10)
        t(ax, cx, by, chip, 15, "white", "bold")
        t(ax, cx, by + 92, head, 17, tc, "bold")
        t(ax, cx, by + 148, lines[0], 13, INK)
        t(ax, cx, by + 182, lines[1], 13, INK)
        if i:
            arrow(ax, x - GAP + 8, by + bh / 2, x - 8, by + bh / 2)

    yb, hb = 470, 86
    card(ax, M + 180, yb, W - 2 * M - 360, hb, fc=PANEL, ec=HAIR, lw=1.4)
    t(ax, W / 2, yb + hb / 2, banner, 16, NAVY, "bold")
    return save(fig, name)


# ── 2) 네 블록 지도 (1600×640 · 2400×960) — W06~W13 강의본 ───────
BLOCKS = [
    ("W1–W4", "기초 · 배분",
     ["자산소유자 · CAPM · SAA", "연기금 모델 · CMA", "MVO · 블랙-리터맨"]),
    ("W5–W9", "포트폴리오 공학",
     ["리스크 패리티 · HRP", "동적 배분 · LDI · GBI", "위험관리 · 성과평가 (FLAM)"]),
    ("W10–W14", "자산군 전략",
     ["주식 퀀트 · 채권 · 팩터", "글로벌 매크로 · CTA", "대체투자 · 비유동성"]),
    ("W15–W16", "통합 · 마무리",
     ["Total Portfolio Approach", "Reference Portfolio", "캡스톤 · 기말시험"]),
]


def blocks(name, title, active, lines_override, banner1, banner2, scale=1.0, star=False):
    W, H = int(1600 * scale), int(640 * scale)
    s = scale
    fig, ax = canvas(W, H)
    t(ax, W / 2, 40 * s, title, 21 * s, NAVY, "bold")

    M, GAP = 68 * s, 44 * s
    n = 4
    bw = (W - 2 * M - (n - 1) * GAP) / n
    by, bh = 138 * s, 280 * s
    for i, (rng, head, lines) in enumerate(BLOCKS):
        lines = lines_override.get(i, lines)
        x = M + i * (bw + GAP)
        if i == 0:
            fc, tcr, tch, tcb = PANEL, MUTED, NAVY, INK
        elif i == 1:
            fc, tcr, tch, tcb = BLUEP, BLUE, BLUE, INK
        elif i == 2:
            fc, tcr, tch, tcb = NAVY, "white", "white", "#e8ecf6"
        else:
            fc, tcr, tch, tcb = GREENP, GREEN, GREEN, INK
        ec, lw = (ORANGE, 3.0 * s) if i == active else (HAIR if i in (0,) else
                                                        (BLUE if i == 1 else
                                                         (NAVY if i == 2 else GREEN)), 1.6 * s)
        card(ax, x, by, bw, bh, fc=fc, ec=ec, lw=lw, r=14 * s)
        cx = x + bw / 2
        t(ax, cx, by + 44 * s, rng, 18 * s, tcr, "bold")
        t(ax, cx, by + 88 * s, head, 18 * s, tch, "bold")
        for j, ln in enumerate(lines[:3]):
            t(ax, cx, by + (146 + 42 * j) * s, ln, 14 * s, tcb)
        if i == active:
            t(ax, cx, by + bh + 34 * s, "지금 여기 ★" if star else "지금 여기",
              15 * s, ORANGE, "bold")

    yb, hb = by + bh + 66 * s, 106 * s
    card(ax, M + 250 * s, yb, W - 2 * M - 500 * s, hb, fc=PANEL, ec=HAIR, lw=1.4)
    t(ax, W / 2, yb + 36 * s, banner1, 16 * s, NAVY, "bold")
    t(ax, W / 2, yb + 72 * s, banner2, 14 * s, MUTED)
    return save(fig, name)


# ── 3) 16주의 여정 (2400×960) — W15 TPA 제1부 마무리 ──────────────
def journey(name="w15_journey.png"):
    W, H = 2400, 960
    fig, ax = canvas(W, H)
    t(ax, W / 2, 58, "16주의 여정 — 모든 도구가 TPA로 수렴한다", 26, NAVY, "bold")

    bars = [
        ("W1–W4 기초 · 배분", "CAPM · SAA · CMA · MVO · 블랙-리터맨",
         "→ Reference Portfolio의 설계 언어", PANEL, HAIR, NAVY),
        ("W5–W9 포트폴리오 공학", "리스크 패리티 · LDI · 위험관리 · FLAM",
         "→ Active Risk Budget과 팩터 렌즈", BLUEP, BLUE, BLUE),
        ("W10–W14 자산군 전략", "퀀트 · 채권 · 팩터 · 매크로 · 대체",
         "→ 위험 예산을 채우는 “아이디어들”", GREENP, GREEN, GREEN),
    ]
    bx, bw, bh, gap, by = 90, 1380, 176, 34, 150
    for i, (head, mid, tail, fc, ec, tc) in enumerate(bars):
        y = by + i * (bh + gap)
        card(ax, bx, y, bw, bh, fc=fc, ec=ec, lw=1.8, r=16)
        t(ax, bx + bw / 2, y + 44, head, 22, tc, "bold")
        t(ax, bx + bw / 2, y + 96, mid, 19, INK)
        t(ax, bx + bw - 40, y + 142, tail, 19, MUTED, ha="right")
        arrow(ax, bx + bw + 14, y + bh / 2, 1600, 446, color="#46506b", lw=2.6)

    nx, ny, nw, nh = 1650, 272, 660, 350
    card(ax, nx, ny, nw, nh, fc=NAVY, ec=NAVY, lw=1.8, r=18)
    t(ax, nx + nw / 2, ny + 62, "Week 15 — TPA", 26, "white", "bold")
    for j, ln in enumerate(["자산군 라벨의 해체",
                            "하나의 기준 · 하나의 위험 예산",
                            "하나의 언어로 경쟁하는 아이디어"]):
        t(ax, nx + nw / 2, ny + 138 + 62 * j, ln, 20, "#e8ecf6")

    yb, hb = 790, 128
    card(ax, 480, yb, 1440, hb, fc=PANEL, ec=HAIR, lw=1.4, r=16)
    t(ax, 1200, yb + 46, "이 과정의 진짜 산출물은 배분표가 아니라 “질문하는 방식”이다", 22, NAVY, "bold")
    t(ax, 1200, yb + 92,
      "— 당신의 부채는 무엇인가 · 기준은 무엇인가 · 그 기준을 이길 아이디어는 무엇인가", 19, MUTED)
    return save(fig, name)


# ── 덱별 정의 ────────────────────────────────────────────────
CHAINS = {
    "w02.png": dict(
        title="커리큘럼에서 오늘의 위치 — 2주",
        steps=[
            ("W1", "자산소유자", ["누구의 돈인가", "4유형 · 거버넌스"], "grey"),
            ("W2 · 오늘", "자산배분과 CAPM", ["BHB 91.5% · SML/CML", "틀린 모델의 쓸모"], "today"),
            ("W3", "CMA · ERP", ["기대수익을 어떻게 정하나", "SAA의 입력값"], "blue"),
            ("W4", "MVO · BL", ["프론티어의 실무화", "뷰의 결합"], "green"),
            ("W12 · W15", "팩터 · TPA", ["어노멀리의 산업화", "벤치마크 설계 · 캡스톤"], "grey"),
        ],
        banner="CAPM의 β는 W3의 ERP, W4의 BL 균형수익률, W12의 팩터, W15의 벤치마크로 계속 산다"),
    "w03.png": dict(
        title="커리큘럼에서 오늘의 위치 — 3주",
        steps=[
            ("W2", "자산배분과 CAPM", ["BHB 91.5% · SML/CML", "틀린 모델의 쓸모"], "grey"),
            ("W3 · 오늘", "연기금 모델과 CMA", ["3대 모델 · μσρ 추정", "MCTR · RC"], "today"),
            ("W4", "MVO · BL", ["CMA를 비중으로 — 추정오차", "뷰의 강건한 결합"], "blue"),
            ("W5", "리스크 패리티 · HRP", ["RC를 같게 — 오늘 MCTR의", "직계 후속"], "green"),
            ("W14 · W15", "대체 · TPA", ["평활화의 귀결 · KIC 케이스의", "최종 답"], "grey"),
        ],
        banner="오늘의 CMA(μσρ)가 W4 MVO의 입력이 되고, 오늘의 RC가 W5 리스크 패리티의 목적함수가 된다"),
    "w04_m4.png": dict(
        title="커리큘럼에서 오늘의 위치 — 4주 · 1교시",
        steps=[
            ("W2", "자산배분과 CAPM", ["BHB 91.5% · SML/CML", "틀린 모델의 쓸모"], "grey"),
            ("W3", "연기금 모델과 CMA", ["3대 모델 · μσρ 추정", "MCTR · RC"], "grey"),
            ("W4 · 오늘", "MVO와 공분산", ["4대 한계 · LW · RMT", "Robust (L1/L2)"], "today"),
            ("W4 · 2교시", "블랙-리터맨", ["균형 + 뷰의 베이즈 결합", "극단 비중의 해소"], "blue"),
            ("W5 · W6", "RP · HRP · Regime", ["역행렬 없이 배분하기", "위기의 공분산"], "green"),
        ],
        banner="W3의 CMA가 오늘 MVO의 입력이었고, 오늘의 한계가 같은 주 BL · W5 HRP · W6 Regime 공분산을 부른다"),
    "w04_m5.png": dict(
        title="커리큘럼에서 오늘의 위치 — 4주 · 2교시",
        steps=[
            ("W3", "연기금 모델과 CMA", ["3대 모델 · μσρ 추정", "MCTR · RC"], "grey"),
            ("W4 · 1교시", "MVO와 공분산", ["4대 한계 · LW · RMT", "Robust (L1/L2)"], "grey"),
            ("W4 · 오늘", "블랙-리터맨", ["역최적화 π · P/Q/Ω", "베이즈 결합"], "today"),
            ("W5", "RP · HRP", ["뷰를 신뢰하지 않는", "정반대의 접근"], "blue"),
            ("W6", "동적 · Regime", ["시간 속의 배분", "위기의 공분산"], "green"),
        ],
        banner="1교시가 “입력을 못 믿겠다”였다면 오늘은 “시장에서 출발하자” — W5는 “아예 μ를 쓰지 말자”로 간다"),
    "w05.png": dict(
        title="커리큘럼에서 오늘의 위치 — 5주",
        steps=[
            ("W4 · 1교시", "MVO와 공분산", ["4대 한계 · LW · RMT", "Robust (L1/L2)"], "grey"),
            ("W4 · 2교시", "블랙-리터맨", ["역최적화 π · P/Q/Ω", "베이즈 결합"], "grey"),
            ("W5 · 오늘", "RP와 HRP", ["μ 없는 배분 · ERC", "역행렬 없는 HRP"], "today"),
            ("W6", "동적 자산배분", ["Merton · 시간 속의 최적", "정적 → 동적"], "blue"),
            ("W9 · W15", "CVaR · TPA", ["꼬리 위험 최적화", "팩터 렌즈의 총화"], "green"),
        ],
        banner="W4 “입력 불신” → “시장에서 출발” → 오늘 “μ를 버린다” — 그리고 W6부터 “시간”이 변수로 들어온다"),
}

MAPS = {
    "w06.png": dict(
        title="16주 로드맵에서 Week 6의 위치 — 정적에서 동적으로 넘어가는 관문",
        active=1,
        lines={1: ["리스크 패리티 · HRP (완료)", "동적 배분 (W6 — 오늘)", "LDI · GBI · 위험관리"]},
        b1="W1–W5의 정적 배분(사진)이 W6에서 동적 배분(영화)으로 — 시간 · 생애주기 · 학습이 들어온다",
        b2="W7 LDI(첫 ★ 마일스톤) · GBI가 이 동적 관점 위에 세워진다"),
    "w07_m8.png": dict(
        title="16주 로드맵에서 Week 7의 위치 — 첫 번째 ★ 마일스톤",
        active=1, star=True,
        lines={1: ["리스크 패리티 · 동적 배분 (완료)", "LDI (W7 — 오늘 ★)", "GBI · 위험관리"]},
        b1="W6의 동적 관점에 “부채”라는 새 좌표축이 더해진다 — 자산이 아니라 자산 대 부채",
        b2="같은 주 GBI(두 번째 ★)는 LDI의 개인 버전 — DB의 기법이 개인 재무로 내려온다"),
    "w07_m9.png": dict(
        title="16주 로드맵에서 Week 7의 위치 — 두 번째 ★ 마일스톤",
        active=1, star=True,
        lines={1: ["리스크 패리티 · 동적 · LDI (완료)", "GBI (W7 — 오늘 ★)", "위험관리 (W9)"]},
        b1="W6 동적 + LDI 부채 관점이 오늘 “개인의 목표”로 통합된다 — LDI의 개인 버전",
        b2="포트폴리오 공학 블록의 마지막 조각은 W9 위험관리 — 만든 것을 지키는 법"),
    "w09.png": dict(
        title="16주 로드맵에서 Week 9의 위치 — 포트폴리오 공학 블록 완결",
        active=1,
        lines={1: ["리스크 패리티 · 동적 · LDI · GBI", "위험관리 · 성과평가 (W9 — 오늘)", "블록 완결"]},
        b1="만드는 법(W5–W7)을 배웠다 — W9는 “지키고 측정하는 법”으로 블록을 닫는다",
        b2="W10부터는 자산군 안으로 — 주식 퀀트 전략에서 다시 만나는 IC와 IR"),
    "w10.png": dict(
        title="16주 로드맵에서 Week 10의 위치 — 도구에서 자산군 전략으로",
        active=2, scale=1.5,
        lines={2: ["주식 퀀트 (W10 — 오늘)", "채권 · 팩터 · 매크로/CTA", "대체투자 · 비유동성"]},
        b1="W9의 Grinold 법칙 IR = IC · √BR 이 W10 퀀트 파이프라인의 수학적 출발점",
        b2="W10 주식 퀀트 → W11 채권 → W12 팩터 투자로 자산군 전략이 이어진다"),
    "w11.png": dict(
        title="16주 로드맵에서 Week 11의 위치 — 자산군 전략의 두 번째 주",
        active=2, scale=1.5,
        lines={2: ["W10 주식 퀀트 (완료)", "채권 (W11 — 오늘)", "팩터 · 매크로/CTA · 대체"]},
        b1="W7 LDI의 듀레이션 매칭이 오늘 볼록성 · 키레이트로 정밀해진다",
        b2="W10 주식 팩터가 오늘 채권 팩터(Carry · Value · Momentum · Defensive)로 확장된다"),
    "w12.png": dict(
        title="16주 로드맵에서 Week 12의 위치 — 자산군 전략의 세 번째 주",
        active=2, scale=1.5,
        lines={2: ["W10 주식 퀀트 · W11 채권 (완료)", "팩터 투자 (W12 — 오늘)", "매크로/CTA · 대체투자"]},
        b1="W10의 FF 모델 · JKP, W11의 채권 팩터가 오늘 하나의 체계로 통합된다",
        b2="W12 팩터 → W13 매크로 · CTA → W14 대체투자로 자산군 전략이 이어진다"),
    "w13.png": dict(
        title="16주 로드맵에서 Week 13의 위치 — 자산군 전략의 네 번째 주",
        active=2, scale=1.5,
        lines={2: ["주식 · 채권 · 팩터 (완료)", "매크로 · CTA (W13 — 오늘)", "다음 주 — 대체투자"]},
        b1="W5 리스크 패리티(All Weather)와 W12 통화 팩터가 오늘 매크로의 언어로 만난다",
        b2="W13 매크로 · CTA → W14 대체투자 → W15 TPA로 과정이 수렴한다"),
}

# 그림이 들어갈 자리 — (pptx 경로, 슬라이드 번호, PNG 이름)
SLOTS = [
    ("W02_자산배분과CAPM/W02_자산배분과CAPM_강의본.pptx", 49, "w02.png"),
    ("W03_연기금모델과CMA/W03_연기금모델과CMA_강의본.pptx", 48, "w03.png"),
    ("W04_MVO와블랙리터맨/W04_M4_MVO와공분산추정_강의본.pptx", 46, "w04_m4.png"),
    ("W04_MVO와블랙리터맨/W04_M5_블랙리터맨_균형과뷰_강의본.pptx", 46, "w04_m5.png"),
    ("W05_리스크패리티와HRP/W05_리스크패리티와HRP_강의본.pptx", 52, "w05.png"),
    ("W06_동적포트폴리오와장기투자/W06_동적포트폴리오와장기투자_강의본.pptx", 52, "w06.png"),
    ("W07_LDI와GBI/W07_M8_LDI_부채연계투자_강의본.pptx", 55, "w07_m8.png"),
    ("W07_LDI와GBI/W07_M9_GBI_목표기반투자_강의본.pptx", 55, "w07_m9.png"),
    ("W09_위험관리와성과평가/W09_위험관리와성과평가_강의본.pptx", 52, "w09.png"),
    ("W10_주식퀀트모델/W10_주식퀀트모델_강의본.pptx", 50, "w10.png"),
    ("W11_채권투자/W11_채권투자_제1부_채권수학과곡선전략.pptx", 53, "w11.png"),
    ("W12_팩터투자/W12_팩터투자_강의본.pptx", 50, "w12.png"),
    ("W13_글로벌매크로와CTA/W13_글로벌매크로와CTA_강의본.pptx", 51, "w13.png"),
    ("W15_TPA/W15_TPA_제1부_사일로해체와설계_강의본.pptx", 58, "w15_journey.png"),
]


def build_all():
    for name, d in CHAINS.items():
        chain(name, d["title"], d["steps"], d["banner"])
    for name, d in MAPS.items():
        blocks(name, d["title"], d["active"], d["lines"], d["b1"], d["b2"],
               scale=d.get("scale", 1.0), star=d.get("star", False))
    journey()


def embed():
    from pptx import Presentation
    root = os.path.dirname(HERE)
    for rel, sn, png in SLOTS:
        path = os.path.join(root, rel)
        prs = Presentation(path)
        pics = [sh for sh in prs.slides[sn - 1].shapes if sh.shape_type == 13]
        if len(pics) != 1:
            print("SKIP (그림이 하나가 아니다)", rel, sn, len(pics)); continue
        pic = pics[0]
        part = pic.part.related_part(pic._element.blip_rId)
        part._blob = open(os.path.join(OUT, png), "rb").read()
        prs.save(path)
        print("embedded", png, "→", rel, f"s{sn}")


if __name__ == "__main__":
    build_all()
    if "--embed" in sys.argv:
        embed()
